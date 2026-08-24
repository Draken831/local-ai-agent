from __future__ import annotations
from dataclasses import dataclass
from pathlib import Path
from typing import Any
import csv
import json

TEXT_EXT = {".txt", ".md", ".log", ".json", ".har", ".csv", ".xml", ".html", ".htm", ".ps1", ".py", ".sh", ".bat", ".cmd", ".yaml", ".yml", ".ini", ".cfg", ".conf", ".sql"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}


@dataclass
class ParsedDocument:
    path: str
    name: str
    extension: str
    size_bytes: int
    parser: str
    text: str
    metadata: dict[str, Any]
    notes: list[str]
    truncated: bool = False


def _read_text(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except Exception:
            pass
    return path.read_text(encoding="utf-8", errors="ignore")


def _trim(text: str, limit: int) -> tuple[str, bool]:
    if len(text) <= limit:
        return text, False
    return text[:limit] + "\n\n[TRUNCATED]", True


def process_document(file_path: str, max_chars: int = 16000, tesseract_cmd: str | None = None) -> ParsedDocument:
    path = Path(file_path).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(path)

    ext = path.suffix.lower()
    parser = "text"
    notes = []
    metadata = {}
    text = ""

    if ext in TEXT_EXT:
        text = _read_text(path)
    elif ext == ".pdf":
        from pypdf import PdfReader
        parser = "pypdf"
        reader = PdfReader(str(path))
        text = "\n\n".join((page.extract_text() or "") for page in reader.pages)
        metadata["pages"] = len(reader.pages)
    elif ext == ".docx":
        from docx import Document
        parser = "python-docx"
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs)
    elif ext == ".xlsx":
        from openpyxl import load_workbook
        parser = "openpyxl"
        book = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in book.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                parts.append("\t".join("" if v is None else str(v) for v in row))
        text = "\n".join(parts)
    elif ext == ".pptx":
        from pptx import Presentation
        parser = "python-pptx"
        deck = Presentation(str(path))
        parts = []
        for i, slide in enumerate(deck.slides, 1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        text = "\n".join(parts)
    elif ext in IMAGE_EXT:
        from PIL import Image
        import pytesseract
        parser = "pytesseract"
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        text = pytesseract.image_to_string(Image.open(path))
        notes.append("OCR output; verify important values against the source image.")
    else:
        raise ValueError(f"Unsupported document type: {ext}")

    text, truncated = _trim(text, max_chars)
    return ParsedDocument(str(path), path.name, ext, path.stat().st_size, parser, text, metadata, notes, truncated)


def document_to_context(file_path: str, question: str = "", max_chars: int = 16000, tesseract_cmd: str | None = None) -> str:
    doc = process_document(file_path, max_chars=max_chars, tesseract_cmd=tesseract_cmd)
    return f"DOCUMENT: {doc.name}\nPARSER: {doc.parser}\nQUESTION: {question}\n\n{doc.text}"
